#!/usr/bin/env python3
"""
PPTX Generator — Generate PowerPoint presentations from JSON spec.

Usage:
    python pptx_generator.py --in spec.json --out output.pptx

Input JSON schema:
{
    "title": "Judul Presentasi",
    "subtitle": "Subjudul (opsional)",
    "slides": [
        {"type": "title"},
        {"type": "bullets", "heading": "...", "bullets": ["..."]},
        {"type": "summary", "heading": "...", "bullets": ["..."], "next_steps": ["..."]}
    ],
    "notes": {
        "enabled": true,
        "per_slide": ["notes slide1", "notes slide2"]
    }
}

@author  Tama El Pablo
@version 1.0.0
"""

import argparse
import json
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# ═══════════════════════════════════════════════════════════
#  STYLE CONSTANTS
# ═══════════════════════════════════════════════════════════

TITLE_FONT_SIZE = Pt(36)
SUBTITLE_FONT_SIZE = Pt(18)
HEADING_FONT_SIZE = Pt(28)
BULLET_FONT_SIZE = Pt(16)
NEXT_STEPS_FONT_SIZE = Pt(14)
NOTES_FONT_SIZE = Pt(12)

COLOR_TITLE = RGBColor(0x1A, 0x23, 0x7E)       # Dark blue
COLOR_SUBTITLE = RGBColor(0x54, 0x54, 0x54)     # Gray
COLOR_HEADING = RGBColor(0x1A, 0x23, 0x7E)      # Dark blue
COLOR_BULLET = RGBColor(0x33, 0x33, 0x33)       # Dark gray
COLOR_NEXT_STEPS = RGBColor(0x0D, 0x74, 0x6E)   # Teal
COLOR_BG_TITLE = RGBColor(0xE8, 0xEA, 0xF6)     # Light blue-gray
COLOR_SUMMARY_BG = RGBColor(0xE0, 0xF2, 0xF1)   # Light teal


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    p.text = notes_text
    p.font.size = NOTES_FONT_SIZE


def build_title_slide(prs, spec, notes_text=None):
    """Build the title/cover slide."""
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)

    set_slide_bg(slide, COLOR_BG_TITLE)

    # Title text box
    left = Inches(1)
    top = Inches(2.2)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = spec.get("title", "Presentasi")
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = spec.get("subtitle", "")
    if subtitle:
        sub_top = Inches(3.8)
        sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(0.8))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = SUBTITLE_FONT_SIZE
        sp.font.color.rgb = COLOR_SUBTITLE
        sp.alignment = PP_ALIGN.CENTER

    if notes_text:
        add_notes(slide, notes_text)

    return slide


def build_bullets_slide(prs, slide_data, notes_text=None):
    """Build a bullet-point content slide."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    heading = slide_data.get("heading", "")
    bullets = slide_data.get("bullets", [])

    # Heading
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(8.4)
    height = Inches(1.0)
    h_box = slide.shapes.add_textbox(left, top, width, height)
    htf = h_box.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hp.text = heading
    hp.font.size = HEADING_FONT_SIZE
    hp.font.bold = True
    hp.font.color.rgb = COLOR_HEADING

    # Divider line
    div_top = Inches(1.4)
    slide.shapes.add_connector(
        1,  # straight
        Inches(0.8), div_top,
        Inches(9.2), div_top,
    )

    # Bullets text box
    b_top = Inches(1.7)
    b_height = Inches(4.5)
    b_box = slide.shapes.add_textbox(left, b_top, width, b_height)
    btf = b_box.text_frame
    btf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = BULLET_FONT_SIZE
        p.font.color.rgb = COLOR_BULLET
        p.space_after = Pt(8)

    if notes_text:
        add_notes(slide, notes_text)

    return slide


def build_summary_slide(prs, slide_data, notes_text=None):
    """Build a summary/takeaways slide."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    set_slide_bg(slide, COLOR_SUMMARY_BG)

    heading = slide_data.get("heading", "Ringkasan")
    bullets = slide_data.get("bullets", [])
    next_steps = slide_data.get("next_steps", [])

    # Heading
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(8.4)
    h_box = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    htf = h_box.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hp.text = heading
    hp.font.size = HEADING_FONT_SIZE
    hp.font.bold = True
    hp.font.color.rgb = COLOR_HEADING

    # Bullets
    b_top = Inches(1.6)
    b_box = slide.shapes.add_textbox(left, b_top, width, Inches(2.5))
    btf = b_box.text_frame
    btf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = f"✓ {bullet}"
        p.font.size = BULLET_FONT_SIZE
        p.font.color.rgb = COLOR_BULLET
        p.space_after = Pt(6)

    # Next Steps section
    if next_steps:
        ns_top = Inches(4.3)
        ns_box = slide.shapes.add_textbox(left, ns_top, width, Inches(2.0))
        nstf = ns_box.text_frame
        nstf.word_wrap = True

        ns_header = nstf.paragraphs[0]
        ns_header.text = "Langkah Selanjutnya:"
        ns_header.font.size = Pt(18)
        ns_header.font.bold = True
        ns_header.font.color.rgb = COLOR_NEXT_STEPS
        ns_header.space_after = Pt(6)

        for step in next_steps:
            p = nstf.add_paragraph()
            p.text = f"→ {step}"
            p.font.size = NEXT_STEPS_FONT_SIZE
            p.font.color.rgb = COLOR_NEXT_STEPS
            p.space_after = Pt(4)

    if notes_text:
        add_notes(slide, notes_text)

    return slide


def generate_pptx(spec, output_path):
    """Generate a PPTX file from a slide spec dict."""
    prs = Presentation()

    # Set 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = spec.get("slides", [])
    notes_cfg = spec.get("notes", {})
    notes_enabled = notes_cfg.get("enabled", False)
    per_slide_notes = notes_cfg.get("per_slide", [])

    for idx, slide_data in enumerate(slides):
        slide_type = slide_data.get("type", "bullets")
        notes_text = per_slide_notes[idx] if notes_enabled and idx < len(per_slide_notes) else None

        if slide_type == "title":
            build_title_slide(prs, spec, notes_text)
        elif slide_type == "summary":
            build_summary_slide(prs, slide_data, notes_text)
        else:
            # Default: bullets
            build_bullets_slide(prs, slide_data, notes_text)

    prs.save(output_path)
    return output_path


def validate_spec(spec):
    """Validate the slide spec schema. Returns list of errors."""
    errors = []

    if not isinstance(spec, dict):
        return ["Spec must be a JSON object"]

    if "title" not in spec or not isinstance(spec["title"], str) or not spec["title"].strip():
        errors.append("'title' is required and must be a non-empty string")

    if "slides" not in spec or not isinstance(spec["slides"], list):
        errors.append("'slides' is required and must be an array")
    else:
        if len(spec["slides"]) < 1:
            errors.append("'slides' must have at least 1 slide")
        if len(spec["slides"]) > 20:
            errors.append("'slides' must have at most 20 slides")

        allowed_types = {"title", "bullets", "summary"}
        for i, s in enumerate(spec["slides"]):
            if not isinstance(s, dict):
                errors.append(f"slides[{i}] must be an object")
                continue
            st = s.get("type", "bullets")
            if st not in allowed_types:
                errors.append(f"slides[{i}].type '{st}' not in {allowed_types}")
            if st == "bullets" and "bullets" in s:
                if not isinstance(s["bullets"], list):
                    errors.append(f"slides[{i}].bullets must be an array")
            if st == "summary" and "bullets" in s:
                if not isinstance(s["bullets"], list):
                    errors.append(f"slides[{i}].bullets must be an array")

    if "notes" in spec:
        notes = spec["notes"]
        if not isinstance(notes, dict):
            errors.append("'notes' must be an object")
        elif "per_slide" in notes and not isinstance(notes["per_slide"], list):
            errors.append("'notes.per_slide' must be an array")

    return errors


def main():
    parser = argparse.ArgumentParser(description="Generate PPTX from JSON spec")
    parser.add_argument("--in", dest="input", required=True, help="Path to input JSON spec file")
    parser.add_argument("--out", dest="output", required=True, help="Path to output PPTX file")
    args = parser.parse_args()

    # Read spec
    try:
        with open(args.input, "r", encoding="utf-8") as f:
            spec = json.load(f)
    except (json.JSONDecodeError, FileNotFoundError) as e:
        print(json.dumps({"success": False, "error": f"Failed to read spec: {e}"}))
        sys.exit(1)

    # Validate
    errors = validate_spec(spec)
    if errors:
        print(json.dumps({"success": False, "errors": errors}))
        sys.exit(1)

    # Generate
    try:
        generate_pptx(spec, args.output)
        file_size = os.path.getsize(args.output)
        print(json.dumps({
            "success": True,
            "output": args.output,
            "slides": len(spec.get("slides", [])),
            "size": file_size,
        }))
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
